#!/usr/bin/env python3
"""
Knowledge File Processor

Polls for knowledge_items with kind='file' and processing_status='ready',
downloads from Supabase Storage, extracts plain text, and marks them
processed so the embedding pipeline can index them.

Supported formats:
  - PDF (.pdf) via pymupdf
  - DOCX (.docx) via python-docx
  - XLSX (.xlsx) via openpyxl
  - CSV (.csv) via stdlib csv
  - PPTX (.pptx) via python-pptx
  - Plain text (.txt, .md, .json, .yaml, etc.)
  - Images (.png, .jpg, etc.) — stored as-is, no text extraction

Environment variables:
  SUPABASE_URL
  SUPABASE_SERVICE_ROLE_KEY

Optional:
  FILE_PROCESS_INTERVAL   — seconds between poll cycles (default: 30)
  FILE_PROCESS_BATCH_SIZE — items per lease (default: 5)
  FILE_PROCESS_LEASE_SEC  — lease duration in seconds (default: 300)

Run:
  python3 /app/file_processor.py
"""

import csv
import io
import json
import os
import time
import traceback
import urllib.parse
import urllib.request
from pathlib import Path

try:
    from registry_config import resolve as resolve_hq_config
except ImportError:
    resolve_hq_config = None  # type: ignore[assignment]

SUPABASE_URL = ""
SUPABASE_KEY = ""
GATEWAY_SLUG = os.environ.get("GATEWAY_ID", "default")
POLL_INTERVAL = int(os.environ.get("FILE_PROCESS_INTERVAL", "30"))
BATCH_SIZE = int(os.environ.get("FILE_PROCESS_BATCH_SIZE", "5"))
LEASE_SECONDS = int(os.environ.get("FILE_PROCESS_LEASE_SEC", "300"))
MAX_TEXT_LENGTH = 500_000


def log(msg: str) -> None:
    print(f"[file_processor] {msg}", flush=True)


def supabase_rpc(fn_name: str, params: dict) -> dict | list | None:
    url = f"{SUPABASE_URL}/rest/v1/rpc/{fn_name}"
    data = json.dumps(params).encode()
    req = urllib.request.Request(
        url,
        data=data,
        headers={
            "apikey": SUPABASE_KEY,
            "Authorization": f"Bearer {SUPABASE_KEY}",
            "Content-Type": "application/json",
            "Prefer": "return=representation",
        },
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=60) as resp:
        body = resp.read().decode()
        return json.loads(body) if body.strip() else None


def download_file(file_url: str) -> bytes:
    if file_url.startswith("http"):
        url = file_url
    else:
        bucket = "assets"
        path = file_url
        url = f"{SUPABASE_URL}/storage/v1/object/{bucket}/{urllib.parse.quote(path)}"

    req = urllib.request.Request(
        url,
        headers={
            "apikey": SUPABASE_KEY,
            "Authorization": f"Bearer {SUPABASE_KEY}",
        },
    )
    with urllib.request.urlopen(req, timeout=120) as resp:
        return resp.read()


def extract_text(data: bytes, mime_type: str | None, file_url: str) -> str | None:
    ext = Path(file_url).suffix.lower() if file_url else ""
    mime = (mime_type or "").lower()

    if mime == "application/pdf" or ext == ".pdf":
        return extract_pdf(data)
    elif mime in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",) or ext == ".docx":
        return extract_docx(data)
    elif mime in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",) or ext == ".xlsx":
        return extract_xlsx(data)
    elif mime == "text/csv" or ext == ".csv":
        return extract_csv(data)
    elif mime in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",) or ext == ".pptx":
        return extract_pptx(data)
    elif mime.startswith("text/") or ext in (
        ".txt",
        ".md",
        ".markdown",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
        ".conf",
        ".log",
        ".sh",
        ".bash",
        ".py",
        ".js",
        ".ts",
        ".tsx",
        ".jsx",
        ".html",
        ".css",
        ".xml",
        ".sql",
        ".env",
        ".gitignore",
    ):
        return data.decode("utf-8", errors="replace")
    elif mime.startswith("image/"):
        return None
    else:
        try:
            return data.decode("utf-8", errors="strict")
        except UnicodeDecodeError:
            return None


def extract_pdf(data: bytes) -> str:
    try:
        import fitz  # pymupdf
    except ImportError:
        log("pymupdf not installed, skipping PDF extraction")
        raise RuntimeError("pymupdf (fitz) not installed — pip install pymupdf")

    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    doc.close()
    return "\n\n".join(pages)


def extract_docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed — pip install python-docx")

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("openpyxl not installed — pip install openpyxl")

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"## {sheet}")
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) if c is not None else "" for c in row]
            if any(v.strip() for v in vals):
                lines.append(", ".join(vals))
    wb.close()
    return "\n".join(lines)


def extract_csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    lines = []
    for row in reader:
        if any(c.strip() for c in row):
            lines.append(", ".join(row))
    return "\n".join(lines)


def extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed — pip install python-pptx")

    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def process_item(item: dict) -> None:
    item_id = item["id"]
    file_url = item.get("file_url", "")
    mime_type = item.get("mime_type")

    log(f"Processing file item (mime={mime_type or 'unknown'})")

    try:
        data = download_file(file_url)
    except Exception as e:
        log(f"  Download failed: {e}")
        supabase_rpc(
            "mark_knowledge_item_processing_failed",
            {
                "p_item_id": item_id,
                "p_error": f"Download failed: {e}",
            },
        )
        return

    try:
        plain_text = extract_text(data, mime_type, file_url)
    except RuntimeError as e:
        log(f"  Extraction error: {e}")
        supabase_rpc(
            "mark_knowledge_item_processing_failed",
            {
                "p_item_id": item_id,
                "p_error": str(e),
            },
        )
        return
    except Exception as e:
        log(f"  Extraction failed: {e}")
        supabase_rpc(
            "mark_knowledge_item_processing_failed",
            {
                "p_item_id": item_id,
                "p_error": f"Text extraction failed: {e}",
            },
        )
        return

    if plain_text is None:
        log("  No text extractable (image or binary), marking done with empty text")
        plain_text = ""

    if len(plain_text) > MAX_TEXT_LENGTH:
        plain_text = plain_text[:MAX_TEXT_LENGTH]

    supabase_rpc(
        "mark_knowledge_item_processed",
        {
            "p_item_id": item_id,
            "p_plain_text": plain_text,
        },
    )
    log(f"  Done, extracted {len(plain_text)} chars")


def poll_cycle() -> int:
    items = supabase_rpc(
        "lease_knowledge_items_for_processing",
        {
            "p_gateway_slug": GATEWAY_SLUG,
            "p_limit": BATCH_SIZE,
            "p_lease_seconds": LEASE_SECONDS,
        },
    )

    if not items:
        return 0

    count = 0
    for item in items:
        try:
            process_item(item)
            count += 1
        except Exception:
            log(f"  Unexpected processing error: {traceback.format_exc()}")
    return count


def resolve_config() -> bool:
    global SUPABASE_URL, SUPABASE_KEY
    SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
    SUPABASE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")

    if not SUPABASE_URL and resolve_hq_config:
        cfg = resolve_hq_config()
        if cfg:
            SUPABASE_URL = cfg.url or ""
            SUPABASE_KEY = cfg.service_role_key or ""

    return bool(SUPABASE_URL and SUPABASE_KEY)


def main() -> None:
    try:
        from sentry_init import init_sentry

        init_sentry("file_processor")
    except ImportError:
        pass

    log("Starting file processor daemon")
    log(f"  poll interval: {POLL_INTERVAL}s, batch: {BATCH_SIZE}, lease: {LEASE_SECONDS}s")

    while not resolve_config():
        log("No Supabase credentials yet, retrying in 10s...")
        time.sleep(10)

    log("Connected to Supabase")

    while True:
        try:
            processed = poll_cycle()
            if processed > 0:
                log(f"Processed {processed} item(s)")
        except Exception as e:
            log(f"Poll cycle error: {traceback.format_exc()}")
            try:
                from sentry_init import capture

                capture(e)
            except ImportError:
                pass

        time.sleep(POLL_INTERVAL)


if __name__ == "__main__":
    main()
